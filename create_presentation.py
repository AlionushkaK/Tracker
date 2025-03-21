import sys
print("Python version:")
print(sys.version)
print("Python executable:")
print(sys.executable)
print("Python path:")
print(sys.path)

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsdecls
import os

def add_hyperlink(shape, url, text, tooltip=None):
    """ Функция для добавления гиперссылки на фигуру """
    hlinkClick = OxmlElement('a:hlinkClick')
    hlinkClick.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', 'rId0')
    if tooltip:
        hlinkClick.set('tooltip', tooltip)

    # Нужно добавить реальный rId, который связывает ссылку с ресурсом
    # Это сложный процесс, требующий генерации и управления реальными relationship внутри пакета pptx
    # Пример с rId0 является упрощенным и не будет работать без реальной реализации управления ссылками

    shape._element.append(hlinkClick)
    shape.text = text

def create_large_presentation(output_file, num_slides=1):
    # Создание объекта презентации
    prs = Presentation()

    for _ in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Использование пустого слайда

        # Создание фигуры для гиперссылки
        left = Inches(1)
        top = Inches(1)
        width = height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        add_hyperlink(shape, "http://example.com/video.mp4", "Смотреть видео", "Кликните для просмотра видео")

    # Сохранение презентации
    prs.save(output_file)
    print(f"Презентация сохранена как {output_file}. Размер файла: {os.path.getsize(output_file) / 1024 ** 2:.2f} MB")

# Путь к файлу презентации
output_file = 'presentation_with_link.pptx'

# Создание презентации
create_large_presentation(output_file)
